"""
Catskill Partners — Template-Based PPTX Generator v3
Clones real slides from both official Catskill decks stored in this directory.
Deck files: api/main_deck.pptx + api/appendix_deck.pptx
"""
from http.server import BaseHTTPRequestHandler
import json, io, os, copy, re, base64, urllib.request
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

_DIR = os.path.dirname(os.path.abspath(__file__))

def _load_decks():
    main_path = os.path.join(_DIR, "main_deck.pptx")
    app_path  = os.path.join(_DIR, "appendix_deck.pptx")
    if not os.path.exists(main_path): raise FileNotFoundError(f"main_deck.pptx not found at {main_path}")
    if not os.path.exists(app_path):  raise FileNotFoundError(f"appendix_deck.pptx not found at {app_path}")
    return Presentation(main_path), Presentation(app_path)

MAIN = {"cover":0,"disclaimer":1,"toc":2,"firm_overview":3,"deal_process":4,"strategy":5,"sector_focus":6,"market_data":7,"pipeline":8,"deal_flow":9,"underwriting":10,"playbook":11,"value_creation":12,"why_now":13,"closing":14,"legal_structure":15,"is_economics":16,"fund_economics":17,"team":18,"deal_summary":19}
APP  = {"appendix_cover":0,"a_legal":1,"a_is_econ":2,"a_fund_econ":3,"a_team":4,"a_pipeline":5,"a_deal_summary":6,"a_sector_adv":7,"a_sector_eng":8,"a_sector_prec":9,"a_cases_cover":10,"a_case_a":11,"a_case_b":12,"a_case_c":13,"a_case_d":14,"a_case_e":15,"a_case_f":16,"a_bios":17}

AUDIENCE_SLIDES = {
    "lp":      [("main","cover"),("main","toc"),("main","firm_overview"),("main","market_data"),("main","strategy"),("main","sector_focus"),("main","pipeline"),("main","value_creation"),("main","why_now"),("app","a_fund_econ"),("app","a_is_econ"),("app","a_team"),("app","a_bios"),("main","closing"),("main","disclaimer")],
    "ib":      [("main","cover"),("main","toc"),("main","firm_overview"),("main","deal_process"),("main","pipeline"),("main","deal_flow"),("main","sector_focus"),("app","a_sector_adv"),("app","a_sector_prec"),("main","value_creation"),("main","why_now"),("app","a_deal_summary"),("app","a_case_a"),("app","a_case_b"),("main","closing"),("main","disclaimer")],
    "founder": [("main","cover"),("main","firm_overview"),("main","strategy"),("main","value_creation"),("main","why_now"),("app","a_case_a"),("app","a_case_b"),("main","closing")],
    "general": [("main","cover"),("main","toc"),("main","firm_overview"),("main","market_data"),("main","strategy"),("main","pipeline"),("main","value_creation"),("main","why_now"),("main","fund_economics"),("app","a_bios"),("main","closing"),("main","disclaimer")],
}

def get_audience_key(a):
    a = a.lower()
    if any(x in a for x in ["lp","institutional","family office","hnw"]): return "lp"
    if any(x in a for x in ["banker","intermediary","ib","broker"]): return "ib"
    if any(x in a for x in ["founder","owner","seller","operator"]): return "founder"
    return "general"

def set_para_text(para, new_text):
    runs = para._p.findall(qn("a:r"))
    if not runs: return
    first_rPr = None
    rPr = runs[0].find(qn("a:rPr"))
    if rPr is not None: first_rPr = copy.deepcopy(rPr)
    for r in runs: para._p.remove(r)
    r_new = etree.SubElement(para._p, qn("a:r"))
    if first_rPr is not None: r_new.append(first_rPr)
    t_new = etree.SubElement(r_new, qn("a:t"))
    t_new.text = new_text

def replace_text(slide, replacements):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            for old, new in replacements.items():
                if old and old in full:
                    set_para_text(para, full.replace(old, new, 1)); break

def clone_slide(src_prs, src_idx, tgt_prs):
    src = src_prs.slides[src_idx]
    blank = tgt_prs.slide_layouts[6]
    dst = tgt_prs.slides.add_slide(blank)
    dst_sp = dst.shapes._spTree
    for child in list(dst_sp): dst_sp.remove(child)
    for elem in src.shapes._spTree: dst_sp.append(copy.deepcopy(elem))
    src_csld = src._element.find(qn("p:cSld"))
    dst_csld = dst._element.find(qn("p:cSld"))
    src_bg = src_csld.find(qn("p:bg"))
    if src_bg is not None:
        dst_bg = dst_csld.find(qn("p:bg"))
        if dst_bg is not None: dst_csld.remove(dst_bg)
        dst_csld.insert(0, copy.deepcopy(src_bg))
    rId_map = {}
    for rId, rel in src.part.rels.items():
        try:
            new_rId = dst.part.relate_to(rel.target if rel.is_external else rel.target_part, rel.reltype, is_external=rel.is_external)
            if new_rId != rId: rId_map[rId] = new_rId
        except: pass
    if rId_map:
        xml_str = etree.tostring(dst._element, encoding="unicode")
        for old, new in rId_map.items():
            xml_str = xml_str.replace(f'r:embed="{old}"', f'r:embed="{new}"')
            xml_str = xml_str.replace(f'r:id="{old}"', f'r:id="{new}"')
        new_elem = etree.fromstring(xml_str)
        parent = dst._element.getparent()
        if parent is not None:
            idx_p = list(parent).index(dst._element); parent.remove(dst._element); parent.insert(idx_p, new_elem)
    return dst

SYSTEM = """You are Morgan Cole, VP of Marketing at Catskill Partners LP.
Generate concise text updates for an investor deck. Return JSON only (no markdown):
{"cover_audience_line":"LP / Institutional Investor Overview — Q1 2026","toc_items":["Firm Overview","Market Opportunity","Investment Strategy","Deal Sourcing","Value Creation","Fund Economics","Team","Appendix"],"custom_text":{}}

FIRM FACTS: Brian Steel CEO/operator >10x MOIC Tenere exit | Mike Fuller 20+ yrs ICT/DC PE | Fund I $250M target | 6-8 platforms $2-20M EBITDA $25-150M EV | 1700+ owner DB | 14 active deals | 25-30% IRR 3-4x MOIC | CLARITY. CRAFT. CAPITAL."""

def call_claude(topic, audience, brief, tone, api_key):
    prompt = f"Deck: {topic} | Audience: {audience} | Brief: {brief or 'Standard overview'} | Tone: {tone}"
    body = json.dumps({"model":"claude-sonnet-4-20250514","max_tokens":800,"system":SYSTEM,"messages":[{"role":"user","content":prompt}]}).encode()
    req = urllib.request.Request("https://api.anthropic.com/v1/messages",data=body,headers={"Content-Type":"application/json","x-api-key":api_key,"anthropic-version":"2023-06-01"})
    try:
        with urllib.request.urlopen(req,timeout=30) as r: d = json.loads(r.read())
        raw = d.get("content",[{}])[0].get("text","{}").strip()
        raw = re.sub(r"```json\n?","",raw).replace("```","").strip()
        return json.loads(raw)
    except: return {"cover_audience_line":audience,"toc_items":[],"custom_text":{}}

def build_deck(topic, audience, brief, tone, task_type, api_key):
    main_prs, app_prs = _load_decks()
    out = Presentation()
    out.slide_width = main_prs.slide_width
    out.slide_height = main_prs.slide_height
    updates = call_claude(topic, audience, brief, tone, api_key)
    audience_key = get_audience_key(audience)
    if task_type in ("lp_update","fund_economics"): audience_key = "lp"
    elif task_type == "ib_teaser": audience_key = "ib"
    slides_to_use = AUDIENCE_SLIDES.get(audience_key, AUDIENCE_SLIDES["general"])
    if task_type == "one_pager":
        slides_to_use = [("main","cover"),("main","firm_overview"),("main","market_data"),("main","why_now"),("main","closing")]
    for src_type, slide_name in slides_to_use:
        src_prs = main_prs if src_type == "main" else app_prs
        idx_map = MAIN if src_type == "main" else APP
        idx = idx_map.get(slide_name)
        if idx is None or idx >= len(src_prs.slides): continue
        slide = clone_slide(src_prs, idx, out)
        replacements = {}
        if slide_name == "cover" and updates.get("cover_audience_line"):
            replacements["Investor Overview"] = updates["cover_audience_line"]
        if slide_name == "toc" and updates.get("toc_items"):
            replacements["Catskill Partners Overview"] = "\n".join(updates["toc_items"][:9])
        if replacements: replace_text(slide, replacements)
    buf = io.BytesIO()
    out.save(buf)
    return buf.getvalue()

class handler(BaseHTTPRequestHandler):
    def log_message(self, *a): pass
    def do_OPTIONS(self): self.send_response(200); self._cors(); self.end_headers()
    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length",0))
            body = json.loads(self.rfile.read(length)) if length else {}
            api_key = os.environ.get("ANTHROPIC_API_KEY","")
            if not api_key: return self._err(500,"API key not configured")
            pptx_bytes = build_deck(body.get("topic","Catskill Partners Overview"),body.get("audience","LP / Institutional Investor"),body.get("brief",""),body.get("tone","Institutional"),body.get("taskType","full_deck"),api_key)
            slug = re.sub(r"[^a-zA-Z0-9]","-",body.get("topic","deck"))[:40]
            filename = f"Catskill-{slug}.pptx"
            resp = json.dumps({"success":True,"filename":filename,"base64":base64.b64encode(pptx_bytes).decode(),"slideCount":len(Presentation(io.BytesIO(pptx_bytes)).slides),"title":body.get("topic","")}).encode()
            self.send_response(200); self._cors()
            self.send_header("Content-Type","application/json"); self.send_header("Content-Length",str(len(resp)))
            self.end_headers(); self.wfile.write(resp)
        except Exception as e:
            import traceback; self._err(500,str(e),traceback.format_exc()[:600])
    def _cors(self):
        self.send_header("Access-Control-Allow-Origin","*")
        self.send_header("Access-Control-Allow-Methods","POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers","Content-Type")
    def _err(self,code,msg,detail=""):
        r=json.dumps({"error":msg,"detail":detail}).encode()
        self.send_response(code); self._cors()
        self.send_header("Content-Type","application/json"); self.send_header("Content-Length",str(len(r)))
        self.end_headers(); self.wfile.write(r)
